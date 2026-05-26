import PyPDF2
import re
from typing import List, Dict
import logging
from app.config import settings

logger = logging.getLogger(__name__)

# Maximum pages/slides/sheet-rows (logical pages) allowed
MAX_PAGES = 500


class DocumentProcessor:
    """
    Extracts text from multiple file formats and splits it into
    page-aware chunks.  Every chunk carries its source page number.
    """

    # ------------------------------------------------------------------
    # PDF  (.pdf)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> List[Dict]:
        """Return list of {text, page_number} per PDF page."""
        pages = []
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                if len(reader.pages) > MAX_PAGES:
                    raise ValueError(
                        f"Document too large. Maximum {MAX_PAGES} pages allowed for accurate results."
                    )
                for idx, page in enumerate(reader.pages, start=1):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append({"text": text, "page_number": idx})
        except ValueError:
            raise
        except Exception as e:
            raise Exception(f"Error reading PDF: {str(e)}")
        return pages

    # ------------------------------------------------------------------
    # DOCX  (.docx)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_docx(file_path: str) -> List[Dict]:
        """Extract paragraphs from a Word document.
        Word docs don't have physical pages at extraction time, so we
        create logical pages of ~3000 chars each."""
        from docx import Document as DocxDocument

        try:
            doc = DocxDocument(file_path)
            full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            return DocumentProcessor._text_to_logical_pages(full_text)
        except Exception as e:
            raise Exception(f"Error reading DOCX: {str(e)}")

    # ------------------------------------------------------------------
    # TXT  (.txt)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_txt(file_path: str) -> List[Dict]:
        """Read plain text file and split into logical pages."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
            return DocumentProcessor._text_to_logical_pages(text)
        except Exception as e:
            raise Exception(f"Error reading TXT: {str(e)}")

    # ------------------------------------------------------------------
    # Markdown  (.md)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_md(file_path: str) -> List[Dict]:
        """Read markdown, strip syntax, split into logical pages."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

            # Strip common markdown syntax
            text = re.sub(r'#{1,6}\s*', '', text)        # headings
            text = re.sub(r'\*\*(.+?)\*\*', r'\1', text) # bold
            text = re.sub(r'\*(.+?)\*', r'\1', text)     # italic
            text = re.sub(r'`{1,3}[^`]*`{1,3}', '', text)  # inline/fenced code
            text = re.sub(r'!\[.*?\]\(.*?\)', '', text)  # images
            text = re.sub(r'\[(.+?)\]\(.*?\)', r'\1', text) # links → keep text
            text = re.sub(r'>\s*', '', text)              # blockquotes
            text = re.sub(r'---+', '', text)              # horizontal rules

            return DocumentProcessor._text_to_logical_pages(text)
        except Exception as e:
            raise Exception(f"Error reading Markdown: {str(e)}")

    # ------------------------------------------------------------------
    # CSV  (.csv)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_csv(file_path: str) -> List[Dict]:
        """Convert CSV rows into readable text chunks, grouped by page."""
        import pandas as pd

        try:
            df = pd.read_csv(file_path)
            if len(df) > MAX_PAGES * 50:  # ~50 rows per logical page
                raise ValueError(
                    f"Document too large. Maximum {MAX_PAGES} pages allowed for accurate results."
                )

            pages = []
            rows_per_page = 50
            columns = list(df.columns)

            for page_start in range(0, len(df), rows_per_page):
                page_num = (page_start // rows_per_page) + 1
                page_rows = df.iloc[page_start:page_start + rows_per_page]
                lines = []
                for _, row in page_rows.iterrows():
                    line_parts = [f"{col}: {row[col]}" for col in columns if pd.notna(row[col])]
                    lines.append(" | ".join(line_parts))
                text = "\n".join(lines)
                if text.strip():
                    pages.append({"text": text, "page_number": page_num})

            if len(pages) > MAX_PAGES:
                raise ValueError(
                    f"Document too large. Maximum {MAX_PAGES} pages allowed for accurate results."
                )
            return pages
        except ValueError:
            raise
        except Exception as e:
            raise Exception(f"Error reading CSV: {str(e)}")

    # ------------------------------------------------------------------
    # PowerPoint  (.pptx)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> List[Dict]:
        """Extract text slide-by-slide. Slide number = page number."""
        from pptx import Presentation

        try:
            prs = Presentation(file_path)
            if len(prs.slides) > MAX_PAGES:
                raise ValueError(
                    f"Document too large. Maximum {MAX_PAGES} pages allowed for accurate results."
                )

            pages = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
                slide_text = "\n".join(texts)
                if slide_text.strip():
                    pages.append({"text": slide_text, "page_number": slide_num})
            return pages
        except ValueError:
            raise
        except Exception as e:
            raise Exception(f"Error reading PPTX: {str(e)}")

    # ------------------------------------------------------------------
    # Excel  (.xlsx)
    # ------------------------------------------------------------------
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> List[Dict]:
        """Extract sheet + row data as text. Each sheet starts a new
        logical page group."""
        import pandas as pd

        try:
            sheets = pd.read_excel(file_path, sheet_name=None, engine='openpyxl')
            pages = []
            page_counter = 0
            rows_per_page = 50

            for sheet_name, df in sheets.items():
                columns = list(df.columns)
                for page_start in range(0, len(df), rows_per_page):
                    page_counter += 1
                    page_rows = df.iloc[page_start:page_start + rows_per_page]
                    lines = [f"[Sheet: {sheet_name}]"]
                    for _, row in page_rows.iterrows():
                        line_parts = [f"{col}: {row[col]}" for col in columns if pd.notna(row[col])]
                        lines.append(" | ".join(line_parts))
                    text = "\n".join(lines)
                    if text.strip():
                        pages.append({"text": text, "page_number": page_counter})

            if page_counter > MAX_PAGES:
                raise ValueError(
                    f"Document too large. Maximum {MAX_PAGES} pages allowed for accurate results."
                )
            return pages
        except ValueError:
            raise
        except Exception as e:
            raise Exception(f"Error reading XLSX: {str(e)}")

    # ------------------------------------------------------------------
    # Chunking  (page-aware)
    # ------------------------------------------------------------------
    @staticmethod
    def chunk_text(text: str, chunk_size: int = None, overlap: int = None) -> List[str]:
        """Split text into chunks with overlap."""
        chunk_size = chunk_size or settings.CHUNK_SIZE
        overlap = overlap or settings.CHUNK_OVERLAP

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + chunk_size
            chunk = text[start:end]

            # Try to break at sentence boundary
            if end < text_length:
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                break_point = max(last_period, last_newline)

                if break_point > chunk_size * 0.5:
                    chunk = chunk[:break_point + 1]
                    end = start + break_point + 1

            chunks.append(chunk.strip())
            start = end - overlap

        return [c for c in chunks if c]

    @staticmethod
    def chunk_pages(pages: List[Dict], chunk_size: int = None, overlap: int = None) -> List[Dict]:
        """
        Takes a list of {text, page_number} and returns a list of
        {text, page_number, chunk_index} where large pages are split
        into multiple chunks but each chunk retains its source page.
        """
        chunk_size = chunk_size or settings.CHUNK_SIZE
        overlap = overlap or settings.CHUNK_OVERLAP
        result = []
        global_chunk_idx = 0
        total_pages = len(pages)

        for i, page in enumerate(pages):
            if i % 50 == 0:
                logger.info(f"Processing page {i+1} of {total_pages}...")
            page_chunks = DocumentProcessor.chunk_text(page["text"], chunk_size, overlap)
            for chunk_text in page_chunks:
                result.append({
                    "text": chunk_text,
                    "page_number": page["page_number"],
                    "chunk_index": global_chunk_idx,
                })
                global_chunk_idx += 1

        return result

    # ------------------------------------------------------------------
    # Dispatcher
    # ------------------------------------------------------------------
    SUPPORTED_EXTENSIONS = ('.pdf', '.docx', '.txt', '.md', '.csv', '.pptx', '.xlsx')

    @staticmethod
    def process_document(file_path: str, filename: str) -> List[Dict]:
        """
        Process a document file and return page-aware chunks.
        Returns: list of {text, page_number, chunk_index}
        """
        ext = filename.lower().rsplit('.', 1)[-1] if '.' in filename else ''
        ext_dot = f'.{ext}'

        extractors = {
            '.pdf':  DocumentProcessor.extract_text_from_pdf,
            '.docx': DocumentProcessor.extract_text_from_docx,
            '.txt':  DocumentProcessor.extract_text_from_txt,
            '.md':   DocumentProcessor.extract_text_from_md,
            '.csv':  DocumentProcessor.extract_text_from_csv,
            '.pptx': DocumentProcessor.extract_text_from_pptx,
            '.xlsx': DocumentProcessor.extract_text_from_xlsx,
        }

        if ext_dot not in extractors:
            raise ValueError(
                "Unsupported file type. Please upload: PDF, DOCX, TXT, MD, CSV, PPTX, or XLSX."
            )

        pages = extractors[ext_dot](file_path)

        if not pages:
            raise ValueError("No text could be extracted from this document.")

        total_pages = len(pages)
        if total_pages > MAX_PAGES:
            raise ValueError(f"Document has {total_pages} pages. Maximum allowed is {MAX_PAGES} pages.")
        elif total_pages > 100:
            logger.warning(f"Large document: {total_pages} pages — processing may be slow")

        return DocumentProcessor.chunk_pages(pages)

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------
    @staticmethod
    def _text_to_logical_pages(text: str, chars_per_page: int = 3000) -> List[Dict]:
        """Split a large flat text into logical pages of roughly
        `chars_per_page` characters, breaking at paragraph boundaries."""
        if not text.strip():
            return []

        pages = []
        current_page = ""
        page_num = 1

        for paragraph in text.split('\n'):
            if len(current_page) + len(paragraph) + 1 > chars_per_page and current_page.strip():
                pages.append({"text": current_page.strip(), "page_number": page_num})
                page_num += 1
                current_page = paragraph + "\n"
            else:
                current_page += paragraph + "\n"

        if current_page.strip():
            pages.append({"text": current_page.strip(), "page_number": page_num})

        if len(pages) > MAX_PAGES:
            raise ValueError(
                f"Document too large. Maximum {MAX_PAGES} pages allowed for accurate results."
            )

        return pages